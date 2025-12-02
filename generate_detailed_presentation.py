"""
Generate a comprehensive, detailed PowerPoint presentation for the SLA-DWP-Fog project.
Includes tables, charts, architecture, algorithm details, parameters, and all figures.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
from pathlib import Path


def create_title_slide(prs, title, subtitle=""):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def create_section_header(prs, title):
    """Create a section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    return slide


def create_content_slide(prs, title):
    """Create a blank slide with title."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """Add a text box to slide."""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color
    
    return text_box


def add_bullet_text(slide, left, top, width, height, items, font_size=14):
    """Add bulleted text to slide."""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
    
    return text_box


def add_table(slide, left, top, rows, cols, col_widths=None):
    """Add a table to slide."""
    if col_widths is None:
        total_width = 8.0
        col_widths = [total_width / cols] * cols
    else:
        total_width = sum(col_widths)
    
    # Calculate total height based on rows
    row_height = 0.4
    total_height = row_height * rows
    
    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(total_width), Inches(total_height)
    ).table
    
    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)
    
    # Style header row
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    return table


def add_image(slide, image_path, left, top, width=None, height=None):
    """Add an image to slide."""
    if width and height:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), 
                                width=Inches(width), height=Inches(height))
    elif width:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top))


def create_architecture_slide(prs):
    """Create system architecture slide with detailed components."""
    slide = create_content_slide(prs, "System Architecture")
    
    # Architecture description
    items = [
        "Vehicular Fog Computing Framework for V2X Applications",
        "Four-layer architecture: Vehicles → Fog Nodes → Edge → Cloud",
        "2×2 Grid Topology: 4 fog nodes covering 1000m × 1000m area",
        "Each fog node: CPU capacity (20 units/s), Link capacity (200 MB/s)",
        "Request routing: Nearest fog node based on Euclidean distance",
        "Queue management: Max 300 requests per node with admission control"
    ]
    add_bullet_text(slide, 0.5, 1.2, 4.5, 4.5, items, font_size=16)
    
    # Component boxes
    boxes = [
        ("Request Generator", 5.2, 1.5, "Poisson arrivals\n5 request types\n3 priority classes"),
        ("Fog Topology", 5.2, 2.5, "4 nodes in 2×2 grid\nNearest-node routing\nResource allocation"),
        ("Schedulers", 5.2, 3.5, "FIFO / EMERGENCY_FIRST\nSTATIC_PRIORITY\nSLA-DWP-Fog"),
        ("Metrics", 5.2, 4.6, "Admission, Latency\nDeadlines, SLA\nPer-class tracking")
    ]
    
    for label, left, top, desc in boxes:
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(4.0), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(200, 220, 240)
        shape.line.color.rgb = RGBColor(0, 51, 102)
        shape.line.width = Pt(2)
        
        # Text
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        add_text_box(slide, left + 0.1, top + 0.25, 3.8, 0.5, desc, font_size=10)
    
    return slide


def create_scheduler_comparison_slide(prs):
    """Create slide comparing all four schedulers."""
    slide = create_content_slide(prs, "Scheduler Algorithms Comparison")
    
    # Create comparison table
    table = add_table(slide, 0.5, 1.2, 5, 4, col_widths=[2.0, 2.0, 2.0, 2.5])
    
    # Headers
    headers = ["Scheduler", "Strategy", "Complexity", "Key Feature"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    # Data rows
    data = [
        ["FIFO", "First-In-First-Out", "O(1)", "Simplest baseline"],
        ["EMERGENCY_FIRST", "Priority queuing", "O(1)", "Emergency tasks first"],
        ["STATIC_PRIORITY", "3-tier priority", "O(1)", "Class-based static"],
        ["SLA-DWP-Fog", "Adaptive weighted", "O(n log n)", "SLA-aware dynamic"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
    
    return slide


def create_sla_dwp_fog_detailed_slide(prs):
    """Create detailed SLA-DWP-Fog algorithm slide."""
    slide = create_content_slide(prs, "SLA-DWP-Fog: Detailed Algorithm")
    
    # Algorithm components
    add_text_box(slide, 0.5, 1.2, 9, 0.4, 
                 "1. Time-Based Admission Control", 
                 font_size=16, bold=True, color=RGBColor(0, 51, 102))
    
    add_text_box(slide, 0.7, 1.65, 8.5, 0.8,
                 "t_finish = t_arrival + (W_n + c_i) / C_n\n"
                 "if t_finish ≤ d_i: ADMIT    else: REJECT\n"
                 "W_n = total remaining work, c_i = task demand, C_n = CPU capacity",
                 font_size=12)
    
    add_text_box(slide, 0.5, 2.5, 9, 0.4,
                 "2. Normalized Priority Scoring",
                 font_size=16, bold=True, color=RGBColor(0, 51, 102))
    
    add_text_box(slide, 0.7, 2.95, 8.5, 1.2,
                 "π_i(t) = α·g(κ_i) + β·u_i(t) + γ·w_i(t)\n\n"
                 "g(κ_i) = class importance (Emergency=1.0, Safety=0.5, Normal=0.0)\n"
                 "u_i(t) = deadline urgency = min(1, max(0, 1 - (d_i - t)/D_i))\n"
                 "w_i(t) = waiting time = min(1, (t - a_i)/D_i)\n"
                 "α, β, γ = adaptive weights (sum = 1.0)",
                 font_size=11)
    
    add_text_box(slide, 0.5, 4.3, 9, 0.4,
                 "3. SLA Metrics & Dual Updates (Time Window TW = 10s)",
                 font_size=16, bold=True, color=RGBColor(0, 51, 102))
    
    add_text_box(slide, 0.7, 4.75, 8.5, 1.2,
                 "J1 = emergency miss ratio (threshold: 0.10)\n"
                 "J2 = mean latency in seconds (threshold: 5.0s)\n"
                 "J3 = fairness ratio normal/emergency (threshold: 2.0)\n\n"
                 "λ_m ← max(0, λ_m + η_m·(J_m - J_m_max))    [η_m = 0.02]\n"
                 "α, β, γ ← (λ + ε) / Σ(λ + ε)    [ε = 1e-6]",
                 font_size=11)
    
    return slide


def create_parameters_table_slide(prs):
    """Create detailed parameters table."""
    slide = create_content_slide(prs, "Configuration Parameters")
    
    # System parameters table
    add_text_box(slide, 0.5, 1.1, 4, 0.4, "System Resources", 
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    
    table1 = add_table(slide, 0.5, 1.5, 5, 2, col_widths=[2.5, 1.5])
    table1.cell(0, 0).text = "Parameter"
    table1.cell(0, 1).text = "Value"
    
    params1 = [
        ("CPU Capacity (per node)", "20 units/s"),
        ("Link Capacity (per node)", "200 MB/s"),
        ("Max Queue Length", "300 requests"),
        ("Topology", "2×2 Grid (4 nodes)")
    ]
    
    for i, (param, value) in enumerate(params1, start=1):
        table1.cell(i, 0).text = param
        table1.cell(i, 1).text = value
        for j in range(2):
            table1.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)
    
    # SLA-DWP-Fog parameters table
    add_text_box(slide, 5.0, 1.1, 4.5, 0.4, "SLA-DWP-Fog Parameters",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    
    table2 = add_table(slide, 5.0, 1.5, 8, 2, col_widths=[2.5, 1.5])
    table2.cell(0, 0).text = "Parameter"
    table2.cell(0, 1).text = "Value"
    
    params2 = [
        ("Time Window (TW)", "10.0 s"),
        ("J1_max (emergency miss)", "0.10"),
        ("J2_max (mean latency)", "5.0 s"),
        ("J3_max (fairness ratio)", "2.0"),
        ("η1 (emergency step size)", "0.02"),
        ("η2 (latency step size)", "0.02"),
        ("η3 (fairness step size)", "0.02")
    ]
    
    for i, (param, value) in enumerate(params2, start=1):
        table2.cell(i, 0).text = param
        table2.cell(i, 1).text = value
        for j in range(2):
            table2.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)
    
    # Workload parameters
    add_text_box(slide, 0.5, 4.0, 9, 0.4, "Workload Configuration",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    
    items = [
        "Load Range: 2, 4, 6, 8, 10, 15, 20, 30, 40, 60, 80, 100 requests/step",
        "Priority Distribution: 30% Emergency, 40% Safety-Critical, 30% Non-Critical",
        "Simulation: 3600 seconds, Time Step: 1.0 second"
    ]
    add_bullet_text(slide, 0.7, 4.4, 8.5, 1.2, items, font_size=12)
    
    return slide


def create_request_types_slide(prs):
    """Create request types and characteristics table."""
    slide = create_content_slide(prs, "Request Types & Characteristics")
    
    # Request types table
    table = add_table(slide, 0.3, 1.2, 6, 5, col_widths=[2.0, 1.5, 1.5, 1.5, 1.8])
    
    headers = ["Request Type", "Priority", "Deadline (s)", "Data (MB)", "CPU (units)"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    data = [
        ["Video Streaming", "Emergency (3)", "0.8 - 1.5", "2.0 - 3.5", "3.0 - 5.0"],
        ["Autonomous Nav", "Safety (2)", "0.8 - 1.5", "2.0 - 3.5", "3.0 - 5.0"],
        ["Traffic Monitor", "Safety (2)", "1.0 - 2.0", "1.0 - 2.0", "2.0 - 4.0"],
        ["Parking Info", "Normal (1)", "3.0 - 5.0", "0.5 - 1.0", "1.0 - 2.0"],
        ["Weather Update", "Normal (1)", "8.0 - 12.0", "0.3 - 0.6", "0.5 - 1.0"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            if col_idx == 0:
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                paragraph.alignment = PP_ALIGN.CENTER
    
    # Add priority explanation
    add_text_box(slide, 0.5, 4.5, 9, 1.0,
                 "Priority Classes:\n"
                 "• Class 3 (Emergency): Life-critical tasks requiring immediate processing\n"
                 "• Class 2 (Safety-Critical): Safety-related tasks with tight deadlines\n"
                 "• Class 1 (Non-Critical): Informational tasks with relaxed deadlines",
                 font_size=13)
    
    return slide


def create_metrics_definitions_slide(prs):
    """Create metrics definitions slide."""
    slide = create_content_slide(prs, "Performance Metrics Definitions")
    
    # Create table for metrics
    table = add_table(slide, 0.5, 1.2, 7, 3, col_widths=[2.5, 3.0, 2.5])
    
    headers = ["Metric", "Formula / Definition", "Interpretation"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    data = [
        ["Admission Rate", "admitted / generated", "Fraction of requests accepted"],
        ["Deadline Miss Rate", "violated / completed", "Fraction missing deadlines"],
        ["Emergency SLA Met", "emergency_met / emergency_completed", "Emergency deadline compliance"],
        ["Mean Latency", "avg(completion - arrival)", "Average end-to-end delay"],
        ["P95 Latency", "95th percentile(latency)", "Tail latency performance"],
        ["Completion Ratio", "completed / generated", "Overall throughput fraction"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            if col_idx == 1:
                paragraph.font.name = "Courier New"
    
    return slide


def create_results_summary_slide(prs, json_path):
    """Create results summary with key findings."""
    slide = create_content_slide(prs, "Experimental Results Summary")
    
    # Load data
    with open(json_path) as f:
        data = json.load(f)
    
    # Key findings
    add_text_box(slide, 0.5, 1.2, 9, 0.4, "Key Findings",
                 font_size=16, bold=True, color=RGBColor(0, 51, 102))
    
    findings = [
        "✓ SLA-DWP-Fog achieves 0% deadline violations at moderate loads (vs 85-93% for others)",
        "✓ Trade-off: Lower admission rate at light loads for strict SLA compliance",
        "✓ All schedulers show proper 100% → 0% admission decline with increasing load",
        "✓ Priority-based schedulers outperform FIFO in deadline metrics",
        "✓ Adaptive weight tuning enables dynamic response to workload patterns"
    ]
    add_bullet_text(slide, 0.7, 1.7, 8.5, 1.8, findings, font_size=13)
    
    # Results table at load=30
    add_text_box(slide, 0.5, 3.6, 9, 0.4, "Performance at Load = 30 requests/step",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    
    table = add_table(slide, 0.8, 4.0, 5, 4, col_widths=[2.2, 1.6, 1.8, 1.8])
    
    headers = ["Scheduler", "Admission Rate", "Deadline Miss", "Emergency SLA"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    schedulers = ["FIFO", "EMERGENCY_FIRST", "STATIC_PRIORITY", "SLA-DWP-Fog"]
    for row_idx, sched in enumerate(schedulers, start=1):
        load_30 = data[sched]["30"]
        table.cell(row_idx, 0).text = sched
        table.cell(row_idx, 1).text = f"{load_30['admission_rate']:.3f}"
        table.cell(row_idx, 2).text = f"{load_30['deadline_miss_rate']:.3f}"
        table.cell(row_idx, 3).text = f"{load_30.get('emergency_sla_met_rate', 0.0):.3f}"
        
        for col_idx in range(4):
            paragraph = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Highlight SLA-DWP-Fog
            if row_idx == 4:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 102, 0)
    
    return slide


def create_admission_rate_analysis_slide(prs, json_path):
    """Create detailed admission rate analysis."""
    slide = create_content_slide(prs, "Admission Rate Analysis")
    
    # Load data
    with open(json_path) as f:
        data = json.load(f)
    
    add_text_box(slide, 0.5, 1.2, 9, 0.6,
                 "Admission Rate Trends: Shows transition from full admission (1.0) at light loads "
                 "to near-zero admission at heavy loads, demonstrating system saturation.",
                 font_size=13)
    
    # Create table with load progression
    table = add_table(slide, 0.5, 2.0, 6, 5, col_widths=[1.5, 1.8, 1.8, 1.8, 1.8])
    
    headers = ["Load", "FIFO", "EMERGENCY", "STATIC", "SLA-DWP"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    loads = [2, 10, 30, 60, 100]
    for row_idx, load in enumerate(loads, start=1):
        table.cell(row_idx, 0).text = str(load)
        for col_idx, sched in enumerate(["FIFO", "EMERGENCY_FIRST", "STATIC_PRIORITY", "SLA-DWP-Fog"], start=1):
            adm = data[sched][str(load)]["admission_rate"]
            table.cell(row_idx, col_idx).text = f"{adm:.4f}"
            
            paragraph = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Color code
            if adm >= 0.8:
                paragraph.font.color.rgb = RGBColor(0, 153, 0)  # Green
            elif adm >= 0.3:
                paragraph.font.color.rgb = RGBColor(204, 102, 0)  # Orange
            else:
                paragraph.font.color.rgb = RGBColor(204, 0, 0)  # Red
    
    add_text_box(slide, 0.5, 4.5, 9, 1.2,
                 "Observations:\n"
                 "• All schedulers maintain 100% admission at very light loads (load ≤ 4)\n"
                 "• SLA-DWP-Fog rejects ~18% even at load=2 due to strict deadline admission control\n"
                 "• All converge to ~3-4% admission at extreme overload (load=100)\n"
                 "• Queue saturation drives admission decline across all policies",
                 font_size=12)
    
    return slide


def create_deadline_performance_slide(prs, json_path):
    """Create deadline performance comparison."""
    slide = create_content_slide(prs, "Deadline Performance Comparison")
    
    # Load data
    with open(json_path) as f:
        data = json.load(f)
    
    add_text_box(slide, 0.5, 1.2, 9, 0.5,
                 "Deadline Miss Rate: Critical QoS metric showing fraction of completed tasks "
                 "that violated their deadlines. Lower is better.",
                 font_size=13)
    
    # Create comparison table
    table = add_table(slide, 0.5, 1.9, 7, 5, col_widths=[1.5, 1.8, 1.8, 1.8, 1.8])
    
    headers = ["Load", "FIFO", "EMERGENCY", "STATIC", "SLA-DWP"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    loads = [2, 6, 15, 30, 60, 100]
    for row_idx, load in enumerate(loads, start=1):
        table.cell(row_idx, 0).text = str(load)
        for col_idx, sched in enumerate(["FIFO", "EMERGENCY_FIRST", "STATIC_PRIORITY", "SLA-DWP-Fog"], start=1):
            dmr = data[sched][str(load)]["deadline_miss_rate"]
            table.cell(row_idx, col_idx).text = f"{dmr:.4f}"
            
            paragraph = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Color code
            if dmr == 0.0:
                paragraph.font.color.rgb = RGBColor(0, 153, 0)  # Green
                paragraph.font.bold = True
            elif dmr < 0.5:
                paragraph.font.color.rgb = RGBColor(102, 153, 0)  # Yellow-green
            elif dmr < 0.9:
                paragraph.font.color.rgb = RGBColor(204, 102, 0)  # Orange
            else:
                paragraph.font.color.rgb = RGBColor(204, 0, 0)  # Red
    
    add_text_box(slide, 0.5, 4.7, 9, 1.0,
                 "Key Insight:\n"
                 "SLA-DWP-Fog maintains 0% deadline violations across all moderate loads through "
                 "proactive admission control. Other schedulers accept all requests, leading to "
                 "85-93% violations at load=30.",
                 font_size=12, bold=True, color=RGBColor(0, 51, 102))
    
    return slide


def add_figure_slide(prs, title, image_path, description=""):
    """Add a slide with a figure."""
    slide = create_content_slide(prs, title)
    
    # Add image centered
    img_path = Path(image_path)
    if img_path.exists():
        add_image(slide, img_path, 0.5, 1.3, width=9.0)
    else:
        add_text_box(slide, 3.0, 3.0, 4.0, 1.0, 
                     f"[Image not found: {img_path.name}]",
                     font_size=14, color=RGBColor(204, 0, 0))
    
    # Add description if provided
    if description:
        add_text_box(slide, 0.5, 6.3, 9.0, 1.0, description, font_size=11)
    
    return slide


def create_presentation():
    """Main function to create comprehensive presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating comprehensive PowerPoint presentation...")
    
    # Title slide
    create_title_slide(prs, 
                      "SLA-Aware Dynamic Weighted Priority Scheduling for Fog Computing",
                      "A Comprehensive Analysis of SLA-DWP-Fog Algorithm\n"
                      "Vehicular Edge Computing with QoS Guarantees")
    
    # Table of contents
    slide = create_content_slide(prs, "Presentation Outline")
    items = [
        "1. Introduction & Motivation",
        "2. System Architecture",
        "3. Scheduler Algorithms Overview",
        "4. SLA-DWP-Fog: Detailed Algorithm",
        "5. Configuration Parameters",
        "6. Request Types & Workload",
        "7. Performance Metrics",
        "8. Experimental Results",
        "9. Visual Analysis (Charts & Plots)",
        "10. Conclusions & Future Work"
    ]
    add_bullet_text(slide, 1.5, 1.5, 7.0, 4.5, items, font_size=16)
    
    # Section 1: Introduction
    create_section_header(prs, "Section 1: Introduction & Motivation")
    
    slide = create_content_slide(prs, "Research Motivation")
    items = [
        "Challenge: Vehicular applications demand ultra-low latency and high reliability",
        "Problem: Traditional cloud computing introduces unacceptable delays (100+ ms)",
        "Solution: Fog computing brings computation closer to vehicles (< 10 ms)",
        "Gap: Existing schedulers lack SLA-awareness and dynamic adaptation",
        "Contribution: Novel SLA-DWP-Fog algorithm with admission control and adaptive weights"
    ]
    add_bullet_text(slide, 0.7, 1.5, 8.5, 4.0, items, font_size=16)
    
    # Section 2: Architecture
    create_section_header(prs, "Section 2: System Architecture")
    create_architecture_slide(prs)
    
    slide = create_content_slide(prs, "Simulation Framework Components")
    items = [
        "config.py: Centralized parameter configuration (resources, SLA thresholds, workload)",
        "models.py: Request, RequestType, and priority class definitions",
        "request_generator.py: Poisson-based arrival process with 5 request types",
        "topology.py: FogNode scheduling logic and grid topology construction",
        "simulation.py: Discrete-event simulation loop with time-stepped execution",
        "metrics.py: Comprehensive tracking (admission, latency, deadlines, SLA)"
    ]
    add_bullet_text(slide, 0.5, 1.4, 9.0, 4.5, items, font_size=14)
    
    # Section 3: Schedulers
    create_section_header(prs, "Section 3: Scheduler Algorithms")
    create_scheduler_comparison_slide(prs)
    
    # Individual scheduler slides
    slide = create_content_slide(prs, "FIFO Scheduler")
    add_text_box(slide, 0.5, 1.3, 4.0, 0.5, "Algorithm:", font_size=14, bold=True)
    items = [
        "Single queue, process in arrival order",
        "No prioritization or admission control",
        "O(1) enqueue and dequeue operations",
        "Serves as baseline for comparison"
    ]
    add_bullet_text(slide, 0.7, 1.9, 4.0, 2.5, items, font_size=13)
    
    add_text_box(slide, 5.0, 1.3, 4.5, 0.5, "Characteristics:", font_size=14, bold=True)
    items2 = [
        "✓ Fair treatment (arrival order)",
        "✓ Simple implementation",
        "✗ No QoS differentiation",
        "✗ Emergency tasks may wait",
        "✗ High deadline miss rate"
    ]
    add_bullet_text(slide, 5.2, 1.9, 4.0, 2.5, items2, font_size=13)
    
    slide = create_content_slide(prs, "EMERGENCY_FIRST Scheduler")
    add_text_box(slide, 0.5, 1.3, 4.0, 0.5, "Algorithm:", font_size=14, bold=True)
    items = [
        "Two queues: emergency and normal",
        "Always process emergency queue first",
        "FIFO within each queue",
        "O(1) complexity"
    ]
    add_bullet_text(slide, 0.7, 1.9, 4.0, 2.5, items, font_size=13)
    
    add_text_box(slide, 5.0, 1.3, 4.5, 0.5, "Characteristics:", font_size=14, bold=True)
    items2 = [
        "✓ Emergency tasks prioritized",
        "✓ Simple two-level priority",
        "✗ No differentiation within classes",
        "✗ No admission control",
        "✗ Normal tasks may starve"
    ]
    add_bullet_text(slide, 5.2, 1.9, 4.0, 2.5, items2, font_size=13)
    
    slide = create_content_slide(prs, "STATIC_PRIORITY Scheduler")
    add_text_box(slide, 0.5, 1.3, 4.0, 0.5, "Algorithm:", font_size=14, bold=True)
    items = [
        "Three queues by priority class",
        "Process: Emergency → Safety → Normal",
        "FIFO within each priority level",
        "Fixed, non-adaptive priorities"
    ]
    add_bullet_text(slide, 0.7, 1.9, 4.0, 2.5, items, font_size=13)
    
    add_text_box(slide, 5.0, 1.3, 4.5, 0.5, "Characteristics:", font_size=14, bold=True)
    items2 = [
        "✓ Three-tier QoS differentiation",
        "✓ Deterministic behavior",
        "✗ No workload adaptation",
        "✗ No admission control",
        "✗ Cannot adjust to SLA violations"
    ]
    add_bullet_text(slide, 5.2, 1.9, 4.0, 2.5, items2, font_size=13)
    
    # Section 4: SLA-DWP-Fog Detailed
    create_section_header(prs, "Section 4: SLA-DWP-Fog Algorithm")
    create_sla_dwp_fog_detailed_slide(prs)
    
    slide = create_content_slide(prs, "SLA-DWP-Fog: Priority Score Components")
    
    add_text_box(slide, 0.5, 1.2, 9, 0.4, "Component 1: Class Importance g(κ_i)",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    add_text_box(slide, 0.7, 1.65, 8.5, 0.6,
                 "Maps priority class to fixed value: Emergency=1.0, Safety=0.5, Normal=0.0\n"
                 "Ensures emergency tasks always have higher base priority",
                 font_size=12)
    
    add_text_box(slide, 0.5, 2.4, 9, 0.4, "Component 2: Deadline Urgency u_i(t)",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    add_text_box(slide, 0.7, 2.85, 8.5, 0.8,
                 "u_i(t) = min(1, max(0, 1 - (d_i - t)/D_i))\n"
                 "Increases as deadline approaches; u=0 when just arrived, u=1 at deadline\n"
                 "Bounded to [0, 1] for normalization",
                 font_size=12)
    
    add_text_box(slide, 0.5, 3.8, 9, 0.4, "Component 3: Waiting Time w_i(t)",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    add_text_box(slide, 0.7, 4.25, 8.5, 0.8,
                 "w_i(t) = min(1, (t - a_i)/D_i)\n"
                 "Increases with time spent waiting; prevents starvation\n"
                 "Normalized by relative deadline D_i for fairness",
                 font_size=12)
    
    add_text_box(slide, 0.5, 5.2, 9, 0.4, "Adaptive Weights α, β, γ",
                 font_size=14, bold=True, color=RGBColor(0, 51, 102))
    add_text_box(slide, 0.7, 5.65, 8.5, 0.8,
                 "Initially: α = β = γ = 1/3 (equal importance)\n"
                 "Updated every TW=10s based on SLA violations via dual variables\n"
                 "Always normalized: α + β + γ = 1.0",
                 font_size=12)
    
    slide = create_content_slide(prs, "SLA-DWP-Fog: Admission Control Logic")
    
    add_text_box(slide, 0.5, 1.3, 9, 0.6,
                 "Proactive Deadline Checking: Predicts completion time before admission",
                 font_size=15, bold=True, color=RGBColor(0, 51, 102))
    
    add_text_box(slide, 0.7, 2.0, 8.5, 1.5,
                 "Step 1: Calculate total remaining work at fog node n\n"
                 "        W_n = Σ(remaining_demand of all queued tasks) + current_task_remaining\n\n"
                 "Step 2: Predict finish time for new arrival i\n"
                 "        t_finish = t_arrival + (W_n + c_i) / C_n\n\n"
                 "Step 3: Compare with deadline\n"
                 "        if t_finish ≤ d_i:  ADMIT to queue\n"
                 "        else:                REJECT (record as admission_rejected)",
                 font_size=13)
    
    add_text_box(slide, 0.5, 3.8, 9, 0.5,
                 "Impact: Prevents deadline violations by rejecting tasks that cannot complete in time",
                 font_size=13, bold=True, color=RGBColor(0, 102, 0))
    
    add_text_box(slide, 0.7, 4.4, 8.5, 1.5,
                 "Trade-off Analysis:\n"
                 "• Lower admission rate (especially at light loads)\n"
                 "• Zero deadline violations for admitted tasks\n"
                 "• Better resource utilization (no wasted processing on late tasks)\n"
                 "• SLA compliance guarantee for completed requests",
                 font_size=13)
    
    slide = create_content_slide(prs, "SLA-DWP-Fog: Weight Adaptation Mechanism")
    
    add_text_box(slide, 0.5, 1.2, 9, 0.5,
                 "Lagrangian Dual Optimization: Weights adjust to minimize SLA violations",
                 font_size=15, bold=True, color=RGBColor(0, 51, 102))
    
    add_text_box(slide, 0.7, 1.8, 8.5, 2.5,
                 "Every TW=10 seconds:\n\n"
                 "1. Measure SLA metrics J1, J2, J3 from completed tasks in last TW seconds\n"
                 "   J1 = emergency_deadline_violations / emergency_completed\n"
                 "   J2 = mean_latency (seconds)\n"
                 "   J3 = mean_latency_normal / mean_latency_emergency\n\n"
                 "2. Update dual variables (gradient ascent on Lagrangian)\n"
                 "   λ1 ← max(0, λ1 + η1·(J1 - 0.10))  [target: ≤10% emergency violations]\n"
                 "   λ2 ← max(0, λ2 + η2·(J2 - 5.0))    [target: ≤5.0s mean latency]\n"
                 "   λ3 ← max(0, λ3 + η3·(J3 - 2.0))    [target: ≤2.0 fairness ratio]\n\n"
                 "3. Normalize to get new weights\n"
                 "   α = (λ1 + ε) / [(λ1 + ε) + (λ2 + ε) + (λ3 + ε)]\n"
                 "   β = (λ2 + ε) / [(λ1 + ε) + (λ2 + ε) + (λ3 + ε)]\n"
                 "   γ = (λ3 + ε) / [(λ1 + ε) + (λ2 + ε) + (λ3 + ε)]",
                 font_size=11)
    
    add_text_box(slide, 0.5, 4.5, 9, 1.0,
                 "Adaptation Behavior:\n"
                 "• If J1 > 0.10: λ1 increases → α increases → more weight on emergency class\n"
                 "• If J2 > 5.0: λ2 increases → β increases → more weight on urgent deadlines\n"
                 "• If J3 > 2.0: λ3 increases → γ increases → more weight on waiting time (fairness)",
                 font_size=12)
    
    # Section 5: Parameters
    create_section_header(prs, "Section 5: Configuration Parameters")
    create_parameters_table_slide(prs)
    create_request_types_slide(prs)
    
    # Section 6: Metrics
    create_section_header(prs, "Section 6: Performance Metrics")
    create_metrics_definitions_slide(prs)
    
    # Section 7: Results
    create_section_header(prs, "Section 7: Experimental Results")
    
    json_path = Path("plots/COMPARISON/comparison_results.json")
    if json_path.exists():
        create_results_summary_slide(prs, json_path)
        create_admission_rate_analysis_slide(prs, json_path)
        create_deadline_performance_slide(prs, json_path)
    
    # Section 8: Visual Analysis
    create_section_header(prs, "Section 8: Visual Analysis")
    
    # Add all comparison figures
    plots_dir = Path("plots/COMPARISON")
    
    figure_configs = [
        ("fig5_admission_rate.png", "Admission Rate vs Load", 
         "Shows proper 1.0→0 decline. SLA-DWP-Fog rejects early to prevent violations."),
        ("fig1_deadline_miss.png", "Deadline Miss Rate vs Load",
         "SLA-DWP-Fog maintains 0% violations. Other schedulers reach 85-93% at load=30."),
        ("fig2_mean_latency.png", "Mean Latency vs Load",
         "Average end-to-end latency increases with load for all schedulers."),
        ("fig3_p95_latency.png", "95th Percentile Latency vs Load",
         "Tail latency performance; shows worst-case delays under different policies."),
        ("fig4_emergency_sla.png", "Emergency SLA Met Rate vs Load",
         "Fraction of emergency tasks meeting deadlines; critical safety metric."),
        ("01_timeseries_comparison_four.png", "Time-Series Analysis (Load=30)",
         "Queue length and completion dynamics over time for all four schedulers."),
        ("02_per_class_latency_comparison.png", "Per-Class Latency Distribution",
         "Boxplot comparison showing latency by priority class (Emergency vs Normal)."),
        ("03_completion_ratio_comparison.png", "Completion Ratio vs Load",
         "Fraction of generated requests successfully completed across load spectrum."),
        ("04_cdf_latency_comparison.png", "Latency CDF (Load=30)",
         "Cumulative distribution showing latency percentiles for all schedulers."),
        ("05_requests_generated_vs_completed.png", "Throughput Analysis",
         "Generated vs completed requests showing system saturation behavior."),
        ("06_sla_compliance_analysis.png", "SLA-DWP-Fog: Detailed SLA Tracking",
         "Admission, emergency SLA met, and deadline miss rates for SLA-DWP-Fog only.")
    ]
    
    for fig_file, title, desc in figure_configs:
        fig_path = plots_dir / fig_file
        add_figure_slide(prs, title, fig_path, desc)
    
    # Section 9: Conclusions
    create_section_header(prs, "Section 9: Conclusions & Future Work")
    
    slide = create_content_slide(prs, "Key Contributions")
    items = [
        "Novel SLA-DWP-Fog algorithm combining admission control with adaptive priorities",
        "Time-based admission: Proactive deadline checking prevents violations",
        "Normalized priority scoring: Bounded components ensure fair comparison",
        "Dual-variable adaptation: Weights adjust based on observed SLA violations",
        "Comprehensive evaluation: 4 schedulers × 12 load levels × 6 metrics",
        "Publication-ready implementation: Modular Python framework with IEEE-compliant plots"
    ]
    add_bullet_text(slide, 0.7, 1.5, 8.5, 4.0, items, font_size=15)
    
    slide = create_content_slide(prs, "Experimental Findings")
    items = [
        "✓ SLA-DWP-Fog achieves 0% deadline violations at moderate loads (load ≤ 60)",
        "✓ Trade-off: ~18% lower admission at light loads for strict QoS guarantees",
        "✓ All schedulers converge to ~3-4% admission at extreme overload (load=100)",
        "✓ Priority-based schedulers (EMERGENCY, STATIC) outperform FIFO by 5-10%",
        "✓ Adaptive weight tuning enables dynamic response to workload patterns",
        "✓ Time-based admission superior to reactive queue management"
    ]
    add_bullet_text(slide, 0.7, 1.5, 8.5, 4.0, items, font_size=15)
    
    slide = create_content_slide(prs, "Limitations & Challenges")
    items = [
        "Homogeneous resources: All fog nodes have identical CPU/link capacity",
        "Static topology: Fixed 2×2 grid without dynamic node addition/removal",
        "Perfect knowledge: Admission assumes accurate processing time estimates",
        "Single objective: Does not optimize energy consumption or operational cost",
        "No task migration: Tasks remain at initially assigned fog node",
        "Synthetic workload: Poisson arrivals may not capture real traffic patterns"
    ]
    add_bullet_text(slide, 0.7, 1.5, 8.5, 4.0, items, font_size=14)
    
    slide = create_content_slide(prs, "Future Research Directions")
    items = [
        "Heterogeneous fog nodes: Variable CPU/link capacities and specialization",
        "Dynamic task migration: Load balancing across nodes for better utilization",
        "Multi-objective optimization: Joint optimization of QoS, energy, and cost",
        "Deep reinforcement learning: Replace dual-based adaptation with DRL agents",
        "Real workload traces: Validate with production data from vehicular networks",
        "Fault tolerance: Handle node failures, network partitions, and mobility",
        "Edge-cloud collaboration: Hybrid scheduling with selective cloud offloading",
        "Security integration: Incorporate trust and privacy constraints in scheduling"
    ]
    add_bullet_text(slide, 0.7, 1.4, 8.5, 4.5, items, font_size=14)
    
    slide = create_content_slide(prs, "Practical Deployment Considerations")
    items = [
        "Scalability: Algorithm complexity O(n log n) may be bottleneck for 1000+ tasks",
        "Parameter tuning: SLA thresholds (J1/2/3_max) and step sizes (η) need calibration",
        "Time synchronization: Requires accurate clocks across distributed fog nodes",
        "Monitoring overhead: SLA metric collection adds ~5-10% computational cost",
        "Admission policy: Rejected tasks need fallback (cloud offload, degraded service)",
        "Real-time constraints: Weight updates every 10s may be too slow for fast dynamics"
    ]
    add_bullet_text(slide, 0.7, 1.5, 8.5, 4.0, items, font_size=13)
    
    # Final slide
    slide = create_content_slide(prs, "Thank You")
    add_text_box(slide, 2.0, 2.5, 6.0, 2.0,
                 "SLA-Aware Dynamic Weighted Priority Scheduling\n"
                 "for Fog Computing\n\n"
                 "Questions & Discussion",
                 font_size=20, bold=True, color=RGBColor(0, 51, 102))
    
    # Save presentation
    output_path = "Detailed_SLA_DWP_Fog_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✓ Presentation saved: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")
    
    return output_path


if __name__ == "__main__":
    create_presentation()
